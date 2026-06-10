#!/usr/bin/env python
"""Generate rich per-week student presentations (Beamer PDF + PPTX).

Each week includes: objectives, the two lectures, an illustration, a
Coincall data-acquisition section (cross-checked against docs.coincall.com),
associative starter Python code, and -- for math-heavy weeks -- derivation
slides that explain WHY each equation is chosen. Week 1 is expanded with the
student-requested progress / timeline / grading / data-sources / resources
content.

One shared content source. PPTX via python-pptx; PDF via LaTeX Beamer.
NC State theme (Wolfpack Red, NCSU wordmark).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

HERE = "/Users/dendisuhubdy/Github/ncsu"
OUT  = os.path.join(HERE, "weekly_presentations")
FIG  = os.path.join(OUT, "figures")
LOGO = os.path.join(HERE, "assets", "ncsu_logo.png")
os.makedirs(OUT, exist_ok=True)

RED, DKRED, GRAY, BLACK, WHITE = (RGBColor(0xCC,0,0), RGBColor(0x99,0,0),
    RGBColor(0x6D,0x6D,0x6D), RGBColor(0x1A,0x1A,0x1A), RGBColor(0xFF,0xFF,0xFF))
TEAL = RGBColor(0x00,0x84,0x73)
FONT = "Arial"
SW, SH = Inches(13.333), Inches(7.5)
MENTOR = "Industry Mentor: Dendi Suhubdy (Bitwyre)  |  Guest: Fenni Kang (Coincall)"

# Coincall API facts (per docs.coincall.com, June 2026). Students must
# re-confirm exact paths against the live docs before each assignment.
COINCALL_BASE = "https://api.coincall.com"

# ============================================================ CONTENT
# Each week is a dict. Optional keys: math (list of (equation, why)),
# data (list of bullet strings), code (caption, code), extra_frames.

def C(s):  # convenience
    return s

WEEKS = []

# ---------------------------------------------------------------- WEEK 1
WEEKS.append(dict(
 stem="week01_introduction", label="Week 1", title="Introduction & Progress Review",
 subtitle="Roadmap, deliverables, grading, data sources, and resources",
 fig="week01.png",
 objectives=[
   "Understand the arc of the program and what you will build",
   "Be clear on the timeline, the three milestones, and how you are graded",
   "Know where the data comes from (Coincall) and how to pull it",
   "Have the Python + PyTorch and C++/pybind11 toolchains running"],
 lec1=("Where We Are (Week 1 Progress Review)", [
   "Maker vs. taker: we build the MAKER -- continuously quoting both sides",
   "The spine of the course is the Avellaneda-Stoikov model",
   "Week 1 status: environment setup, data access, and this orientation",
   "By Week 5 you ship a calibrated quoting engine (Milestone 1)"]),
 lec2=("Tooling & How to Get Help", [
   "Problem sets: Python + PyTorch (autograd does differentiation for you)",
   "Fast computation (Week 7): C++ exposed to Python via pybind11",
   "Mentorship is asynchronous (email/Telegram); tutors run day-to-day",
   "Office hours with tutors; come with a minimal reproducible snippet"]),
 reading=["Cartea-Jaikumar-Penalva, Ch. 1; Avellaneda-Stoikov (2008) -- skim",
          "Coincall API docs: https://docs.coincall.com/  (read 'Getting Started')"],
 # Week 1 gets several extra, student-requested frames defined below.
 extra=[
   ("Timeline & Deliverable Clarification", [
     "10 weeks. Week 1 introduction (done). Weeks 2-10 technical.",
     "Milestone 1 (Wk 5): Avellaneda-Stoikov engine, Python/PyTorch",
     "Milestone 2 (Wk 7): fast pricing/Greeks library, C++ via pybind11",
     "Milestone 3 (Wk 9): full options market-making backtest",
     "Final (Wk 10): one advanced extension + report (<=30pp) + 20-min talk",
     "Problem sets are due Sunday 23:59 ET; milestones Friday of the week"]),
   ("Syllabus & Grading Walkthrough", [
     "Problem sets 30%  |  Milestones (3) 30%  |  Final 30%  |  Participation 10%",
     "Code graded on correctness, reproducibility, quality -- in that order",
     "Numerical-verification problems: right setup + right error magnitude",
     "Milestone 2 latency is a HARD GATE (sub-5ms surface) -- caps at 70 if missed",
     "Late PS: -20% until Tue; nothing after. Milestones: no late submission.",
     "Set random seeds: reproducibility is explicitly graded"]),
   ("Data Sources Discussion (Coincall)", [
     "Exchange: Coincall.  REST base: " + COINCALL_BASE,
     "Provided as Parquet: BTC/ETH options L2 + trades, perp L2, funding, SVI",
     "Coincall orderbook is L2 (aggregated depth) -- there is NO L3 feed",
     "You can also pull live via the REST/WebSocket API (next slides)",
     "Auth: API key + HMAC-SHA256 signature; market data is largely public",
     "Always cross-check exact endpoints at https://docs.coincall.com/"]),
   ("Additional Resources -- Math-Heavy Sections", [
     "Stochastic control: Pham (2009), Ch. 3-4 (HJB, verification)",
     "Market making theory: Gueant (2016), Ch. 4; Cartea et al., Ch. 10-11",
     "Options/vol surface: Gatheral (2006); Gatheral-Jacquier (2014)",
     "Hedging: Taleb (1997), gamma scalping & the gamma-theta identity",
     "El Aoud-Abergel (2015): multi-Greek option market making (Week 8)",
     "We DERIVE the key results in lecture and explain every modeling choice"]),
 ],
 data=[
   "Base URL: " + COINCALL_BASE + "   |   Python SDK: github.com/CoincallExchange",
   "Signed request headers: X-CC-APIKEY, sign (HMAC-SHA256, upper hex), ts (ms), X-REQ-TS-DIFF",
   "Signing: sort params alphabetically, append uuid, ts, x-req-ts-diff, HMAC-SHA256 with secret",
   "Rate limits (per docs): orderbook 30/2s per user; public config 10/2s per IP",
   "Key endpoints you will use across the program:",
   "  Instruments:   GET /open/option/getInstruments/{baseCurrency}",
   "  Option chain:  GET /open/option/get/v1/{index}",
   "  Option book:   GET /open/option/order/orderbook/v1/{symbol}  (100 depth)",
   "  Option detail: GET /open/option/detail/v1/{symbol}  (mark, IV, delta/gamma/vega/theta)",
   "  Trades:        GET /open/option/trade/lasttrade/v1/{symbol}",
   "  Futures book:  GET /open/futures/order/orderbook/v1/{symbol}",
   "  Funding:       GET /open/public/fundingRate/v1?symbol=BTCUSD",
   "  Klines:        GET /open/option/market/kline/history/v1/{optionName}?period=h1",
   "WebSocket channels: /options/pricing, /options/orderbook, /options/lasttrade, /options/kline",
   "ALWAYS verify the exact path/params against https://docs.coincall.com/ -- the API evolves."],
 code=("coincall_client.py -- signed REST helper (verify signing vs. live docs)", r'''
import time, hmac, hashlib, uuid, requests

BASE = "https://api.coincall.com"

def signed_get(path, params=None, api_key="", api_secret=""):
    """Minimal Coincall signed GET. Public market-data endpoints often work
    unsigned; this shows the signed pattern used for account endpoints."""
    params = params or {}
    ts, window = str(int(time.time() * 1000)), "5000"
    items = sorted(params.items())                       # sort alphabetically
    payload = "&".join(f"{k}={v}" for k, v in items)
    payload += f"&uuid={uuid.uuid4().hex}&ts={ts}&x-req-ts-diff={window}"
    sign = hmac.new(api_secret.encode(), payload.encode(),
                    hashlib.sha256).hexdigest().upper()
    headers = {"X-CC-APIKEY": api_key, "sign": sign,
               "ts": ts, "X-REQ-TS-DIFF": window}
    r = requests.get(BASE + path, params=params, headers=headers, timeout=10)
    r.raise_for_status()
    return r.json()

# Example: list all BTC option instruments
print(signed_get("/open/option/getInstruments/BTC"))
'''),
))

# ---------------------------------------------------------------- WEEK 2
WEEKS.append(dict(
 stem="week02_microstructure", label="Week 2", title="Foundations of Market Microstructure",
 subtitle="From the limit order book to the role of the market maker",
 fig="week02.png",
 objectives=[
   "Reconstruct an L2 order book from Coincall data",
   "Decompose the bid-ask spread (Roll, Glosten-Milgrom, Kyle)",
   "Explain the market maker as a continuous liquidity provider"],
 lec1=("Lecture 1: Order Book Mechanics", [
   "Order types: market, limit, IOC, post-only, hidden",
   "Price-time priority and the matching algorithm",
   "L1 and L2 views (Coincall provides L2; there is no L3 / per-order feed)",
   "Because data is L2, we model aggregate intensity, not queue position"]),
 lec2=("Lecture 2: Spread Decomposition", [
   "Adverse selection + inventory + order-processing costs",
   "Roll (1984): effective spread from trade prices",
   "Glosten-Milgrom (1985): asymmetric information sets the spread",
   "Realized spread, effective spread, price impact"]),
 reading=["Cartea-Jaikumar-Penalva, Ch. 1-3; Roll (1984); Glosten-Milgrom (1985)"],
 data=[
   "Goal: pull a single option's order book and recent trades, build mid/spread.",
   "Option order book (100 depth):  GET /open/option/order/orderbook/v1/{symbol}",
   "Recent trades:                   GET /open/option/trade/lasttrade/v1/{symbol}",
   "List symbols first:              GET /open/option/getInstruments/BTC",
   "Response gives arrays of [price, size] for bids and asks -> aggregate (L2)",
   "For replay/backtests use the provided Parquet L2 snapshots (1s for options)",
   "Live updates: subscribe to the /options/orderbook WebSocket channel"],
 code=("week2_orderbook.py -- mid & spread from a Coincall L2 book", r'''
from coincall_client import signed_get   # from Week 1

def l2_top_of_book(symbol, key, secret):
    ob = signed_get(f"/open/option/order/orderbook/v1/{symbol}", {}, key, secret)
    data = ob["data"]                       # confirm field names in the docs
    bids = sorted(data["bids"], key=lambda x: -float(x[0]))
    asks = sorted(data["asks"], key=lambda x:  float(x[0]))
    best_bid, best_ask = float(bids[0][0]), float(asks[0][0])
    mid    = 0.5 * (best_bid + best_ask)
    spread = best_ask - best_bid
    return mid, spread, best_bid, best_ask

# For an offline backtest, replay the provided Parquet instead:
#   import pandas as pd
#   book = pd.read_parquet("coincall_btc_options_l2_2025Q4.parquet")
#   book["mid"] = 0.5 * (book.bid_px_0 + book.ask_px_0)
'''),
))

# ---------------------------------------------------------------- WEEK 3
WEEKS.append(dict(
 stem="week03_avellaneda_stoikov", label="Week 3", title="The Avellaneda-Stoikov Model",
 subtitle="Optimal market making via stochastic control -- derived from scratch",
 fig="week03.png",
 objectives=[
   "Set up the Avellaneda-Stoikov control problem",
   "Follow the HJB derivation and understand every modeling choice",
   "Interpret the reservation price and the optimal half-spread"],
 lec1=("Lecture 1: Setup & Dynamics", [
   "Mid-price: dS_t = sigma dW_t  (arithmetic Brownian motion)",
   "Fill intensity: lambda(delta) = A exp(-kappa delta)",
   "Wealth from fills; inventory q jumps by +/-1 on a fill",
   "Exponential (CARA) utility of terminal wealth"]),
 lec2=("Lecture 2: HJB, Ansatz, Closed Form", [
   "Dynamic programming -> HJB partial differential equation",
   "Separation ansatz collapses HJB to ODEs in inventory",
   "Asymptotic (long-horizon) closed-form quotes",
   "Reservation price skews quotes against inventory"]),
 math=[
   (r"dS_t = \sigma\, dW_t",
    r"Arithmetic (not geometric) Brownian motion is chosen because over a short "
    r"market-making horizon price changes are small and roughly additive; it keeps "
    r"the HJB linear in $S$ and makes the cash dynamics exact."),
   (r"\lambda(\delta) = A\, e^{-\kappa \delta}",
    r"Fill intensity decays with quote distance $\delta$. The exponential form is "
    r"the simplest law consistent with 'further from mid $\Rightarrow$ fewer fills'; "
    r"$A$ sets the base rate, $\kappa$ the decay. It also yields a closed-form optimum."),
   (r"\max\ \mathbb{E}\!\left[-e^{-\gamma\,(X_T + q_T S_T)}\right]",
    r"CARA (exponential) utility is used because constant absolute risk aversion makes "
    r"the value function factor out wealth, so the optimal quotes do not depend on the "
    r"cash level -- exactly what we want for a market maker."),
   (r"\partial_t V + \tfrac12\sigma^2 \partial_{SS} V "
    r"+ \max_{\delta^b}\lambda(\delta^b)\big[V(q{+}1,x{-}S{+}\delta^b)-V\big]"
    r"+ \max_{\delta^a}\lambda(\delta^a)\big[V(q{-}1,x{+}S{+}\delta^a)-V\big]=0",
    r"The HJB equation: the diffusion term prices inventory risk; each $\max$ term is "
    r"the dealer optimally choosing the bid/ask offset, trading fill probability against "
    r"captured spread."),
   (r"V(t,x,S,q) = -e^{-\gamma(x+qS)}\,\theta(t,q)",
    r"This ansatz is not a guess from nowhere: under CARA $+$ arithmetic dynamics the "
    r"value function inherits a multiplicative structure across $(x, qS)$, so all the "
    r"wealth/price dependence is exponential and only $\theta(t,q)$ remains to solve."),
   (r"r(t,q) = S - q\,\gamma\,\sigma^2 (T-t)",
    r"The reservation (indifference) price: it sits BELOW mid when you are long "
    r"($q>0$) so you quote to sell down inventory. The skew scales with risk aversion "
    r"$\gamma$, variance $\sigma^2$, and time left -- more risk $\Rightarrow$ more skew."),
   (r"\delta^* = \tfrac{1}{\gamma}\ln\!\big(1+\tfrac{\gamma}{\kappa}\big) "
    r"+ \tfrac{\gamma\sigma^2}{2}(T-t)",
    r"The optimal half-spread = an order-flow term $\tfrac1\gamma\ln(1+\gamma/\kappa)$ "
    r"(how hard fills are to get) plus an inventory-risk term $\tfrac{\gamma\sigma^2}{2}(T-t)$. "
    r"You widen when risk-averse, when volatile, and when far from the horizon."),
 ],
 reading=["Avellaneda-Stoikov (2008); Cartea-Jaikumar-Penalva, Ch. 10"],
 data=["This week uses simulated dynamics (no live data needed).",
       "For realistic sigma, estimate it from Coincall perp klines:",
       "  GET /open/option/market/kline/history/v1/{name}?period=m1  (or futures klines)",
       "You will calibrate A, kappa from real fills in Week 5."],
 code=("week3_as_quotes.py -- closed-form quotes in PyTorch", r'''
import torch

def as_quotes(S, q, gamma, sigma, kappa, tau):
    """Asymptotic Avellaneda-Stoikov bid/ask. Vectorized over inventory q."""
    reservation = S - q * gamma * sigma**2 * tau
    half_spread = (1.0/gamma) * torch.log1p(gamma/kappa) \
                  + 0.5 * gamma * sigma**2 * tau
    return reservation - half_spread, reservation + half_spread  # bid, ask

S = torch.tensor(100.0); q = torch.arange(-10, 11).float()
bid, ask = as_quotes(S, q, gamma=0.1, sigma=2.0, kappa=1.5, tau=1.0)
print("inventory ->", q.tolist())
print("bid       ->", bid.round(decimals=2).tolist())
'''),
))

# ---------------------------------------------------------------- WEEK 4
WEEKS.append(dict(
 stem="week04_control", label="Week 4", title="Stochastic Optimal Control",
 subtitle="The general theory behind the Avellaneda-Stoikov worked example",
 fig="week04.png",
 objectives=[
   "Apply the dynamic programming principle in continuous time",
   "Understand the HJB equation and the verification theorem",
   "See why CARA / CRRA utilities give tractable problems"],
 lec1=("Lecture 1: Dynamic Programming & HJB", [
   "Controlled diffusions and admissible controls",
   "The dynamic programming principle (Bellman)",
   "The infinitesimal generator and the HJB equation"]),
 lec2=("Lecture 2: Verification & Utilities", [
   "Verification theorem: a smooth HJB solution IS the value function",
   "CARA vs CRRA: where wealth enters the optimal policy",
   "Why market making uses CARA"]),
 math=[
   (r"V(t,x)=\sup_{u}\ \mathbb{E}\!\left[\int_t^T f(s,X_s,u_s)\,ds + g(X_T)\right]",
    r"The value function is the best achievable expected reward from $(t,x)$ onward. "
    r"We define the whole problem through $V$ because Bellman's principle then gives a "
    r"local (differential) equation instead of an intractable global optimization."),
   (r"V(t,x)=\sup_u \mathbb{E}\big[V(t{+}dt, X_{t+dt}) + f\,dt\big]",
    r"The dynamic programming principle: an optimal policy is optimal from every later "
    r"state too. Expanding $V(t+dt,\cdot)$ with Ito's lemma turns this into the HJB PDE."),
   (r"\partial_t V + \sup_{u}\big\{\mathcal{L}^u V + f(t,x,u)\big\}=0,\quad V(T,x)=g(x)",
    r"The HJB equation. $\mathcal{L}^u$ is the generator of the controlled diffusion; the "
    r"$\sup$ is the instantaneous optimal control. This is the engine behind every model "
    r"in the course -- Avellaneda-Stoikov is one special case."),
   (r"\text{CARA: } U(w)=-e^{-\gamma w}\ \Rightarrow\ \text{policy independent of } w",
    r"Constant absolute risk aversion makes marginal utility scale-free in wealth, so the "
    r"optimal quotes do not move with the cash balance. CRRA ($U=w^{1-\gamma}/(1-\gamma)$) "
    r"instead makes the policy scale WITH wealth -- natural for investment, awkward for a desk."),
 ],
 reading=["Pham (2009), Ch. 3-4; Cartea-Jaikumar-Penalva, Ch. 6"],
 data=["Theory week -- no Coincall data required.",
       "Optional: estimate drift/vol from perp klines to parameterize Merton."],
 code=("week4_merton.py -- HJB value iteration (CARA) in PyTorch", r'''
import torch
# Merton under CARA, solved by backward value iteration on a wealth grid.
w = torch.linspace(-5, 5, 401); dt = 1e-3; mu, r, sig, gam = 0.08, 0.02, 0.2, 1.0
V = -torch.exp(-gam * w)                      # terminal utility
for _ in range(1000):                         # march backward in time
    # optimal risky holding pi* = (mu - r)/(gam*sig^2) is wealth-independent (CARA)
    pi = (mu - r) / (gam * sig**2)
    drift = (r*w + pi*(mu - r))
    diff  = 0.5 * (pi*sig)**2
    dV = torch.gradient(V, spacing=(w[1]-w[0]).item())[0]
    d2V = torch.gradient(dV, spacing=(w[1]-w[0]).item())[0]
    V = V + dt * (drift*dV + diff*d2V)
print("pi* (constant in wealth):", (mu - r)/(gam*sig**2))
'''),
))

# ---------------------------------------------------------------- WEEK 5
WEEKS.append(dict(
 stem="week05_engine_milestone1", label="Week 5", title="The Avellaneda-Stoikov Engine",
 subtitle="Implementation, calibration, and backtest -- Milestone 1",
 fig="week05.png",
 objectives=[
   "Implement the quotes as a tested PyTorch engine",
   "Calibrate the fill intensities by maximum likelihood",
   "Backtest on one month of Coincall BTC-PERP data"],
 lec1=("Lecture 1: From Closed Form to Code", [
   "Engine architecture; state from the L2 book",
   "Map continuous-time quantities to discrete events",
   "Calibrate lambda(delta) = A exp(-kappa delta) from fills"]),
 lec2=("Lecture 2: Calibration & Validation", [
   "Maximum-likelihood fit of A and kappa with goodness-of-fit",
   "Replay engine + fill simulator + P&L attribution",
   "Reproducibility under a fixed seed"]),
 math=[
   (r"\mathcal{L}(A,\kappa)=\prod_i \lambda(\delta_i)\,e^{-\lambda(\delta_i)\,\Delta t_i}",
    r"Fills are a (doubly stochastic) Poisson process, so the likelihood of the observed "
    r"fills given quote distances $\delta_i$ is a product of Poisson densities. We fit by "
    r"MAXIMUM LIKELIHOOD rather than least-squares because the data are event counts, not "
    r"Gaussian residuals."),
   (r"-\ln\mathcal{L} = \sum_i\big[\lambda(\delta_i)\Delta t_i - \ln\lambda(\delta_i)\big]",
    r"Minimizing this negative log-likelihood is convex in $\ln A$ and well-behaved in "
    r"$\kappa$; PyTorch autograd gives the gradients, so a few hundred Adam steps converge."),
 ],
 reading=["Avellaneda-Stoikov (2008) [implementation]; Gueant (2016), Ch. 4"],
 data=[
   "Build the intensity-calibration set from trades vs. your quotes:",
   "  Trades:  GET /open/option/trade/lasttrade/v1/{symbol}  (or trade history)",
   "  Book:    GET /open/futures/order/orderbook/v1/{symbol}  for the mid",
   "Offline: coincall_btc_perp_trades_2025Q4.parquet + ..._l2_2025Q4.parquet",
   "A pre-aggregated coincall_intensity_calibration_set.parquet is provided",
   "delta_i = |fill price - mid at fill|; group fills by delta bucket for the MLE"],
 code=("week5_calibrate.py -- intensity MLE with PyTorch autograd", r'''
import torch
# delta: quote distances at which fills occurred; dt: exposure time per fill
delta = torch.tensor([...]); dt = torch.tensor([...])
logA  = torch.zeros(1, requires_grad=True)
kappa = torch.ones(1,  requires_grad=True)
opt = torch.optim.Adam([logA, kappa], lr=0.05)
for step in range(2000):
    lam = torch.exp(logA) * torch.exp(-kappa.clamp(min=1e-3) * delta)
    nll = (lam * dt - torch.log(lam)).sum()          # Poisson negative log-lik
    opt.zero_grad(); nll.backward(); opt.step()
print("A =", torch.exp(logA).item(), " kappa =", kappa.item())
'''),
))

# ---------------------------------------------------------------- WEEK 6
WEEKS.append(dict(
 stem="week06_vol_surface", label="Week 6", title="Options Pricing & the Volatility Surface",
 subtitle="Greeks, no-arbitrage, and SVI/SABR calibration",
 fig="week06.png",
 objectives=[
   "Recall the Greeks, including vanna and volga",
   "State the no-arbitrage constraints on the surface",
   "Calibrate SVI and SABR to Coincall data"],
 lec1=("Lecture 1: Greeks & the Surface", [
   "Black-Scholes and the Greek vector",
   "Vanna, volga, and volatility risk",
   "Calendar and butterfly arbitrage constraints"]),
 lec2=("Lecture 2: SVI & SABR Calibration", [
   "Raw and arbitrage-free SVI",
   "SVI by nonlinear least squares (autograd Jacobians)",
   "The SABR / Hagan expansion"]),
 math=[
   (r"w(k)=a+b\big(\rho(k-m)+\sqrt{(k-m)^2+\varsigma^2}\big)",
    r"Raw SVI for total implied variance $w=\sigma^2 T$ in log-moneyness $k$. The "
    r"hyperbola shape is chosen because it matches observed smiles with only 5 parameters "
    r"and has closed-form slopes for the no-arbitrage checks."),
   (r"\text{butterfly: } g(k)\ge 0,\qquad \text{calendar: } \partial_T w(k,T)\ge 0",
    r"These inequalities are exactly 'no free option-portfolio arbitrage': a negative "
    r"$g(k)$ implies a negative risk-neutral density, and decreasing total variance in $T$ "
    r"implies a calendar arbitrage. We enforce them so the fitted surface is tradeable."),
 ],
 reading=["Gatheral (2006), Ch. 1-5; Gatheral-Jacquier (2014); Hagan et al. (2002)"],
 data=[
   "Pull a full option chain with mark IV and Greeks per strike:",
   "  Chain:   GET /open/option/get/v1/{index}        (index = BTC/ETH underlying)",
   "  Detail:  GET /open/option/detail/v1/{symbol}    (mark, IV, delta/gamma/vega/theta/rho)",
   "  Symbols: GET /open/option/getInstruments/BTC",
   "Offline: coincall_btc_options_chain_snapshots_2025Q4.parquet (5-min snapshots)",
   "A baseline coincall_calibrated_svi_surfaces_2025Q4.parquet is provided"],
 code=("week6_bs_autograd.py -- Greeks via PyTorch autograd (no hand AD)", r'''
import torch
def bs_call(S, K, T, r, sig):
    d1 = (torch.log(S/K) + (r + 0.5*sig**2)*T) / (sig*torch.sqrt(T))
    d2 = d1 - sig*torch.sqrt(T)
    N = torch.distributions.Normal(0., 1.).cdf
    return S*N(d1) - K*torch.exp(-r*T)*N(d2)

S,K,T,r,sig = [torch.tensor(x, requires_grad=True)
               for x in (100., 100., 1.0, 0.0, 0.6)]   # crypto IV ~ 60%
price = bs_call(S,K,T,r,sig)
vega,  = torch.autograd.grad(price, sig, create_graph=True)
vanna, = torch.autograd.grad(vega, S, create_graph=True)   # d vega / dS
volga, = torch.autograd.grad(vega, sig)                    # d vega / d sig
print(float(price), float(vega), float(vanna), float(volga))
'''),
))

# ---------------------------------------------------------------- WEEK 7
WEEKS.append(dict(
 stem="week07_fast_computation_milestone2", label="Week 7",
 title="Fast Computation: Prices, Greeks, Implied Vol, Surfaces",
 subtitle="The latency budget, in C++ via pybind11 -- Milestone 2",
 fig="week07.png",
 objectives=[
   "Implement COS / Carr-Madan pricers and robust implied-vol in C++",
   "Expose the C++ kernels to Python with pybind11",
   "Validate Greeks against PyTorch autograd; hit the 5 ms budget"],
 lec1=("Lecture 1: CF Pricing & Implied Vol", [
   "Characteristic functions: Black-Scholes and Heston",
   "Fang-Oosterlee COS method (the workhorse)",
   "Robust implied-vol inversion (good guess + safeguarded Newton)"]),
 lec2=("Lecture 2: Binding C++ with pybind11", [
   "Structure a pybind11 module; build with CMake",
   "Pass NumPy/PyTorch tensors with no copy (buffer protocol)",
   "Release the GIL around the C++ hot loop",
   "PyTorch autograd is the REFERENCE for Greeks; C++ is the speed"]),
 math=[
   (r"C(K)=e^{-rT}\,\tfrac{1}{\pi}\!\int_0^\infty \mathrm{Re}\!\big[e^{-iu\ln K}\hat{\phi}(u)\big]\,du",
    r"Once we leave Black-Scholes (e.g. Heston), there is no closed-form price, but the "
    r"characteristic function $\phi$ IS known. Pricing becomes a single Fourier integral -- "
    r"this is why CF methods exist and why they are fast."),
   (r"C \approx \sum_{n=0}^{N-1}{}' \mathrm{Re}\!\Big[\phi\!\big(\tfrac{n\pi}{b-a}\big)e^{-i n\pi\frac{a}{b-a}}\Big]V_n",
    r"The Fang-Oosterlee COS method replaces the integral with a cosine series that "
    r"converges geometrically for smooth densities. We truncate $[a,b]$ from the cumulants, "
    r"not by guessing -- that is what makes it reach $10^{-4}$ accuracy in microseconds."),
 ],
 reading=["Carr-Madan (1999); Fang-Oosterlee (2008); Jaeckel (2015); pybind11 docs"],
 data=[
   "Inputs are the calibrated surface from Week 6 plus reference prices:",
   "  coincall_btc_options_chain_snapshots_2025Q4.parquet",
   "  coincall_calibrated_svi_surfaces_2025Q4.parquet (baseline)",
   "  coincall_reference_prices_2025Q4.parquet (accuracy target, 1e-4)",
   "Live mark/IV/Greeks for spot checks: GET /open/option/detail/v1/{symbol}"],
 code=("fastmm pybind11 module -- C++ kernel exposed to Python", r'''
// fastmm.cpp  (build with CMake + pybind11; import as `fastmm`)
#include <pybind11/pybind11.h>
#include <pybind11/numpy.h>
namespace py = pybind11;

double cos_price(double S, double K, double T, double r, double sigma);  // COS pricer

py::array_t<double> price_surface(py::array_t<double> K, double S,
                                  double T, double r, double sigma) {
    auto k = K.unchecked<1>(); auto out = py::array_t<double>(k.shape(0));
    auto o = out.mutable_unchecked<1>();
    { py::gil_scoped_release nogil;                 // release GIL for the hot loop
      for (py::ssize_t i = 0; i < k.shape(0); ++i)
          o(i) = cos_price(S, k(i), T, r, sigma); }
    return out;                                     // zero-copy back to NumPy
}
PYBIND11_MODULE(fastmm, m) { m.def("price_surface", &price_surface); }
'''),
))

# ---------------------------------------------------------------- WEEK 8
WEEKS.append(dict(
 stem="week08_options_mm_hedging", label="Week 8", title="Options Market Making & Dynamic Hedging",
 subtitle="Avellaneda-Stoikov with a vector-valued (Greek) inventory",
 fig="week08.png",
 objectives=[
   "Set up the El Aoud-Abergel multi-Greek framework",
   "Understand the quadratic inventory penalty",
   "Quantify the P&L of a delta-hedged book"],
 lec1=("Lecture 1: Multi-Greek Inventory & HJB", [
   "Inventory becomes a vector q = (Delta, Gamma, vega, ...)",
   "Quadratic inventory penalty; reduces to AS in 1-D",
   "Multi-Greek HJB and the separation ansatz"]),
 lec2=("Lecture 2: Dynamic Hedging", [
   "Continuous vs discrete delta hedging",
   "The gamma-theta-variance identity",
   "Perp futures as the hedge instrument; funding cost"]),
 math=[
   (r"\text{penalty} = \tfrac12\, q^{\mathsf T} \Sigma\, q",
    r"Inventory risk for an options book is multi-dimensional (delta, gamma, vega...). "
    r"A quadratic form is the local (mean-variance) risk of the Greek vector; $\Sigma$ is "
    r"calibrated to the empirical covariance of the Greeks, so cross-risks are priced, not "
    r"just diagonal ones. It reduces to AS's $q^2$ penalty when there is one Greek."),
   (r"d\Pi = \tfrac12\,\Gamma\,S^2\big(\sigma_{\text{real}}^2-\sigma_{\text{impl}}^2\big)\,dt",
    r"The gamma-theta-variance identity: a delta-hedged option earns (or pays) the gap "
    r"between realized and implied variance, scaled by dollar gamma. This single equation "
    r"explains ~80% of a hedged book's daily P&L and motivates why we track gamma at all."),
 ],
 reading=["El Aoud-Abergel (2015); Taleb (1997); Cartea et al., Ch. 11"],
 data=[
   "Greeks per contract (for the inventory vector and Sigma):",
   "  GET /open/option/detail/v1/{symbol}  -> delta, gamma, vega, theta, rho",
   "Hedge instrument + carry:",
   "  Perp book:  GET /open/futures/order/orderbook/v1/{symbol}",
   "  Funding:    GET /open/public/fundingRate/v1?symbol=BTCUSD",
   "Offline: coincall_funding_rates_2025Q4.parquet for hedge-cost accounting"],
 code=("week8_hedge.py -- aggregate Greeks & threshold delta hedge", r'''
import torch
def book_greeks(positions, greeks):
    # positions: (n,), greeks: (n,5) columns [delta,gamma,vega,theta,rho]
    return (positions[:, None] * greeks).sum(0)      # aggregate inventory vector

def delta_hedge(net_delta, perp_pos, threshold=0.1):
    if net_delta.abs() > threshold:                  # rebalance only when breached
        trade = -net_delta                           # hedge with perp futures
        return perp_pos + trade, trade
    return perp_pos, torch.tensor(0.0)
'''),
))

# ---------------------------------------------------------------- WEEK 9
WEEKS.append(dict(
 stem="week09_backtesting_taker_milestone3", label="Week 9",
 title="Backtesting, Risk Analysis & Taker Strategies",
 subtitle="Assembling the full system -- Milestone 3 (Coincall guest session)",
 fig="week09.png",
 objectives=[
   "Build an event-driven options backtest",
   "Produce a full P&L attribution and risk metrics",
   "Understand taker strategies and order-flow toxicity"],
 lec1=("Lecture 1: Backtest Architecture & Risk", [
   "Event-driven simulation with deterministic replay",
   "Fill simulator; state of inventory and Greeks",
   "Pitfalls: look-ahead, unrealistic fills, expiry handling",
   "Sharpe with block-bootstrap confidence intervals"]),
 lec2=("Lecture 2 (Guest: Fenni Kang, Coincall): Taker Strategies", [
   "When and why takers cross the spread",
   "Signals behind aggressive flow: momentum, liquidations, news",
   "What taker behaviour means for the maker's toxicity"]),
 math=[
   (r"\widehat{SR}=\frac{\bar r}{s_r}\sqrt{252},\quad "
    r"\text{CI via block bootstrap on } \{r_t\}",
    r"Daily Sharpe is annualized by $\sqrt{252}$, but a point estimate over a few months "
    r"is noisy and serially correlated. We resample BLOCKS of returns (preserving "
    r"autocorrelation) to get an honest confidence interval -- a high Sharpe with a wide CI "
    r"is not a real edge."),
 ],
 reading=["Cartea et al., Ch. 10-11; Bailey-Lopez de Prado (2014); Lopez de Prado (2018)"],
 data=[
   "Full options L2 + trades + perp + funding for 3 months:",
   "  coincall_btc_options_l2_2025Q4.parquet / ..._trades_...",
   "  coincall_eth_options_l2_2025Q4.parquet / ..._trades_...",
   "  coincall_btc_eth_perp_l2_2025Q4.parquet ; coincall_funding_rates_...",
   "Klines for regime context: GET /open/option/market/kline/history/v1/{name}?period=h1",
   "Trade tape carries mark IV at trade -> calibrate the fill simulator"],
 code=("week9_attribution.py -- P&L attribution + bootstrap Sharpe", r'''
import numpy as np
def attribute(spread, inventory, gamma, hedging):
    total = spread + inventory + gamma + hedging
    assert abs(total.sum() - (spread.sum()+inventory.sum()
               +gamma.sum()+hedging.sum())) < 1e-6           # must sum to total
    return dict(spread=spread, inventory=inventory, gamma=gamma,
                hedging=hedging, total=total)

def bootstrap_sharpe(r, block=5, n=2000, seed=0):
    rng = np.random.default_rng(seed); T = len(r); out = []
    for _ in range(n):
        idx = (rng.integers(0, T-block, T//block)[:,None] + np.arange(block)).ravel()
        s = r[idx]; out.append(s.mean()/s.std()*np.sqrt(252))
    return np.percentile(out, [2.5, 50, 97.5])               # CI for the Sharpe
'''),
))

# ---------------------------------------------------------------- WEEK 10
WEEKS.append(dict(
 stem="week10_advanced_structured_products", label="Week 10",
 title="Advanced Avellaneda-Stoikov, Structured Products & Finals",
 subtitle="Extensions, the corporate view, and communicating results",
 fig="week10.png",
 objectives=[
   "Implement one advanced Avellaneda-Stoikov extension",
   "Understand corporate structured products",
   "Write and present a research-grade report"],
 lec1=("Lecture 1: Advanced Extensions", [
   "Multi-asset market making (BTC + ETH, correlated)",
   "Adverse-selection-aware and signal-driven quoting",
   "Robust optimization: worst-case inventory penalty"]),
 lec2=("Lecture 2 (Guest: Fenni Kang, Coincall): Structured Products", [
   "Collars to fix a price band",
   "Accumulators / decumulators for staged flow",
   "Coupon / yield-enhancement products",
   "How that flow lands on the market-maker's book"]),
 math=[
   (r"\min_{\delta}\ \max_{\Sigma\in\mathcal{U}}\ \Big\{\text{spread P\&L} - \tfrac12 q^{\mathsf T}\Sigma q\Big\}",
    r"Robust market making: we do not trust a single covariance estimate $\Sigma$, so we "
    r"optimize against the WORST case over an uncertainty set $\mathcal{U}$. This widens "
    r"quotes exactly where parameter risk is highest -- a principled cure for over-fitting "
    r"the penalty matrix."),
 ],
 reading=["Gueant (2016) [chosen extension]; El Aoud-Abergel (2015); +1 paper of your choice"],
 data=[
   "Multi-asset extension needs BTC and ETH chains together:",
   "  GET /open/option/get/v1/BTC   and   GET /open/option/get/v1/ETH",
   "Correlation from perp klines (BTC vs ETH) over your backtest window",
   "Reuse all Milestone-3 Parquet datasets; no new sources required"],
 code=("week10_collar.py -- structured-product payoff (collar)", r'''
import numpy as np
def collar_payoff(S, spot=100.0, put_K=90.0, call_K=110.0):
    """Long underlying + long put(put_K) + short call(call_K): a fixed band."""
    return (S - spot) + np.maximum(put_K - S, 0.0) - np.maximum(S - call_K, 0.0)

S = np.linspace(60, 140, 9)
print(dict(zip(S.round(0), collar_payoff(S).round(1))))   # capped above, floored below
'''),
))


# ============================================================ GLOSSARIES
# ASCII-only (the weekly Beamer decks compile without a Unicode header).
GLOSSARIES = {
 "week01_introduction": [
   ("Market maker", "liquidity provider that continuously quotes bid and ask, earning the spread"),
   ("Taker", "participant that crosses the spread for immediate execution"),
   ("Bid-ask spread", "best ask minus best bid; the maker's revenue and risk premium"),
   ("Inventory", "the maker's net position; the main source of risk"),
   ("Milestone", "graded software deliverable, due in Weeks 5, 7, and 9"),
   ("PyTorch / autograd", "Python ML library; automatic differentiation computes Greeks for you"),
   ("pybind11", "library that exposes C++ functions to Python (used in Week 7)")],
 "week02_microstructure": [
   ("Limit order book (LOB)", "all resting limit orders organized by price level"),
   ("L1 / L2 / L3", "top-of-book / aggregated depth / per-order; Coincall provides L2 only"),
   ("Mid-price", "the average of the best bid and best ask"),
   ("Adverse selection", "loss from trading against better-informed counterparties"),
   ("Effective spread", "actual cost relative to mid (estimated by the Roll measure)"),
   ("Price impact", "how far the mid moves in the direction of an incoming trade"),
   ("Order-flow toxicity", "the degree to which incoming flow is informed")],
 "week03_avellaneda_stoikov": [
   ("Reservation price", "inventory-adjusted indifference price that quotes are centered on"),
   ("Half-spread", "distance from the reservation price to each quote"),
   ("HJB equation", "the PDE characterizing the optimal value function"),
   ("Ansatz", "an assumed functional form that reduces the HJB to ODEs"),
   ("CARA utility", "constant absolute risk aversion, U(w) = -exp(-gamma w)"),
   ("Intensity lambda(delta)", "Poisson fill rate as a function of quote distance delta"),
   ("Risk aversion gamma", "parameter setting inventory skew and spread width")],
 "week04_control": [
   ("Stochastic optimal control", "choosing controls to optimize an expected objective under randomness"),
   ("Dynamic programming principle", "an optimal policy stays optimal from every later state"),
   ("Value function", "the best achievable expected reward from a given state onward"),
   ("Infinitesimal generator", "operator giving the expected instantaneous change of a function"),
   ("Verification theorem", "conditions under which an HJB solution is the value function"),
   ("CRRA utility", "constant relative risk aversion; the policy scales with wealth"),
   ("Admissible control", "a control satisfying the problem's constraints")],
 "week05_engine_milestone1": [
   ("Maximum likelihood (MLE)", "fitting parameters by maximizing the likelihood of the data"),
   ("Poisson process", "counting process for random arrivals (here, fills)"),
   ("Goodness-of-fit (chi-square)", "a statistic measuring how well the fit matches the data"),
   ("Backtest", "simulation of a strategy on historical data"),
   ("Fill simulator", "model of which of the maker's quotes get executed"),
   ("P&L attribution", "decomposition of profit into its sources"),
   ("Reproducibility", "identical results on re-run under a fixed random seed")],
 "week06_vol_surface": [
   ("Implied volatility", "the volatility that matches the market option price under Black-Scholes"),
   ("Volatility smile", "implied volatility as a function of strike / moneyness"),
   ("SVI", "Stochastic Volatility Inspired parameterization of total variance"),
   ("SABR", "stochastic alpha-beta-rho model with the Hagan implied-vol expansion"),
   ("Greeks", "price sensitivities: delta, gamma, vega, theta, rho"),
   ("Vanna / Volga", "second-order Greeks: d vega / dS and d vega / d sigma"),
   ("Calendar / butterfly arbitrage", "no-arbitrage constraints the surface must satisfy")],
 "week07_fast_computation_milestone2": [
   ("Characteristic function", "the Fourier transform of the log-return density"),
   ("COS method", "Fourier-cosine series option pricing (Fang-Oosterlee)"),
   ("Carr-Madan", "FFT-based option pricing via a damped payoff"),
   ("Implied-vol inversion", "solving for the IV that reproduces a price (safeguarded Newton)"),
   ("pybind11", "header-only C++ <-> Python binding"),
   ("GIL", "Python's Global Interpreter Lock; released around the C++ hot loop"),
   ("Latency budget", "the maximum allowed time (sub-5 ms full-surface revaluation)")],
 "week08_options_mm_hedging": [
   ("Greek inventory vector", "aggregated (delta, gamma, vega, ...) exposure of the book"),
   ("Quadratic penalty", "the risk term (1/2) q^T Sigma q on the Greek vector"),
   ("Delta hedging", "neutralizing directional risk using the underlying or a perpetual"),
   ("Gamma scalping", "P&L earned by re-hedging a long-gamma position"),
   ("Gamma-theta-variance identity", "hedged P&L = (1/2) Gamma S^2 (realized var - implied var) dt"),
   ("Funding rate", "periodic payment between perpetual-swap longs and shorts"),
   ("El Aoud-Abergel", "the multi-Greek option market-making framework")],
 "week09_backtesting_taker_milestone3": [
   ("Event-driven backtest", "a simulation that replays timestamped events in order"),
   ("Deterministic replay", "bitwise-identical results on re-run"),
   ("Look-ahead bias", "using information not available at decision time"),
   ("Sharpe ratio", "risk-adjusted return: annualized mean over standard deviation"),
   ("Block bootstrap", "resampling blocks of returns for a CI that keeps autocorrelation"),
   ("Maximum drawdown", "the largest peak-to-trough decline in equity"),
   ("Taker strategy", "aggressive, spread-crossing execution")],
 "week10_advanced_structured_products": [
   ("Multi-asset market making", "jointly quoting correlated underlyings (e.g., BTC and ETH)"),
   ("Robust optimization", "optimizing against the worst case over an uncertainty set"),
   ("Adverse-selection-aware quoting", "skewing quotes using order-flow signals"),
   ("Collar", "long put + short call that fixes a price band around a holding"),
   ("Accumulator / decumulator", "products for staged accumulation or liquidation"),
   ("Coupon product", "a yield-enhancement structured note"),
   ("Structured product", "a packaged derivative payoff sold to clients or corporates")],
}


# ============================================================ REFERENCES
# Full citations per week (ASCII; the weekly Beamer decks have no Unicode header).
REFERENCES = {
 "week01_introduction": [
   "Avellaneda, M., & Stoikov, S. (2008). High-frequency trading in a limit order book. Quantitative Finance, 8(3), 217-224.",
   "Gueant, O. (2016). The Financial Mathematics of Market Liquidity. Chapman & Hall/CRC.",
   "Cartea, A., Jaikumar, S., & Penalva, J. (2015). Algorithmic and High-Frequency Trading. Cambridge Univ. Press.",
   "Coincall API documentation. https://docs.coincall.com/"],
 "week02_microstructure": [
   "Roll, R. (1984). A simple implicit measure of the effective bid-ask spread. Journal of Finance, 39(4), 1127-1139.",
   "Glosten, L. R., & Milgrom, P. R. (1985). Bid, ask and transaction prices in a specialist market. J. Financial Economics, 14(1), 71-100.",
   "Kyle, A. S. (1985). Continuous auctions and insider trading. Econometrica, 53(6), 1315-1335.",
   "O'Hara, M. (1995). Market Microstructure Theory. Blackwell.",
   "Cartea, Jaikumar, & Penalva (2015). Algorithmic and High-Frequency Trading, Ch. 1-3."],
 "week03_avellaneda_stoikov": [
   "Avellaneda, M., & Stoikov, S. (2008). High-frequency trading in a limit order book. Quantitative Finance, 8(3), 217-224.",
   "Gueant, O. (2016). The Financial Mathematics of Market Liquidity, Ch. 4. Chapman & Hall/CRC.",
   "Cartea, Jaikumar, & Penalva (2015). Algorithmic and High-Frequency Trading, Ch. 10."],
 "week04_control": [
   "Pham, H. (2009). Continuous-time Stochastic Control and Optimization with Financial Applications. Springer.",
   "Fleming, W. H., & Soner, H. M. (2006). Controlled Markov Processes and Viscosity Solutions (2nd ed.). Springer.",
   "Oksendal, B. (2003). Stochastic Differential Equations (6th ed.). Springer.",
   "Almgren, R., & Chriss, N. (2001). Optimal execution of portfolio transactions. Journal of Risk, 3(2), 5-40."],
 "week05_engine_milestone1": [
   "Avellaneda, M., & Stoikov, S. (2008). High-frequency trading in a limit order book. Quantitative Finance, 8(3), 217-224.",
   "Gueant, O. (2016). The Financial Mathematics of Market Liquidity, Ch. 4. Chapman & Hall/CRC.",
   "Cartea, Jaikumar, & Penalva (2015). Algorithmic and High-Frequency Trading, Ch. 10.",
   "Coincall API documentation. https://docs.coincall.com/"],
 "week06_vol_surface": [
   "Gatheral, J. (2006). The Volatility Surface: A Practitioner's Guide. Wiley.",
   "Gatheral, J., & Jacquier, A. (2014). Arbitrage-free SVI volatility surfaces. Quantitative Finance, 14(1), 59-71.",
   "Hagan, P. S., Kumar, D., Lesniewski, A. S., & Woodward, D. E. (2002). Managing smile risk. Wilmott Magazine, Sept., 84-108.",
   "Hull, J. C. (2021). Options, Futures, and Other Derivatives (11th ed.). Pearson."],
 "week07_fast_computation_milestone2": [
   "Carr, P., & Madan, D. (1999). Option valuation using the fast Fourier transform. J. Computational Finance, 2(4), 61-73.",
   "Fang, F., & Oosterlee, C. W. (2008). A novel pricing method for European options based on Fourier-cosine series expansions. SIAM J. Sci. Comput., 31(2), 826-848.",
   "Jackel, P. (2015). Let's be rational. Wilmott Magazine.",
   "Jakob, W., Rhinelander, J., & Moldovan, D. (2017). pybind11. https://github.com/pybind/pybind11",
   "Paszke, A., et al. (2019). PyTorch: an imperative style, high-performance deep learning library. NeurIPS."],
 "week08_options_mm_hedging": [
   "El Aoud, S., & Abergel, F. (2015). A stochastic control approach to option market making. Market Microstructure and Liquidity, 1(1), 1550006.",
   "Taleb, N. N. (1997). Dynamic Hedging: Managing Vanilla and Exotic Options. Wiley.",
   "Cartea, Jaikumar, & Penalva (2015). Algorithmic and High-Frequency Trading, Ch. 11."],
 "week09_backtesting_taker_milestone3": [
   "Cartea, Jaikumar, & Penalva (2015). Algorithmic and High-Frequency Trading, Ch. 10-11.",
   "Bailey, D. H., & Lopez de Prado, M. (2014). The deflated Sharpe ratio. Journal of Portfolio Management, 40(5), 94-107.",
   "Lopez de Prado, M. (2018). Advances in Financial Machine Learning. Wiley.",
   "Coincall API documentation. https://docs.coincall.com/"],
 "week10_advanced_structured_products": [
   "Gueant, O. (2016). The Financial Mathematics of Market Liquidity. Chapman & Hall/CRC.",
   "El Aoud, S., & Abergel, F. (2015). A stochastic control approach to option market making. Market Microstructure and Liquidity, 1(1), 1550006.",
   "Cartea, Jaikumar, & Penalva (2015). Algorithmic and High-Frequency Trading.",
   "Plus one paper of your choice relevant to the chosen extension."],
}


# ============================================================ PPTX RENDER
def _set(run, size, color, bold=False, italic=False, mono=False):
    run.font.name = "Consolas" if mono else FONT
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic


def _title_slide(prs, w):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    band = s.shapes.add_shape(1, 0, 0, SW, Inches(2.5))
    band.fill.solid(); band.fill.fore_color.rgb = RED; band.line.fill.background()
    s.shapes.add_picture(LOGO, Inches(0.5), Inches(0.5), width=Inches(3.0))
    tb = s.shapes.add_textbox(Inches(0.6), Inches(2.9), Inches(12.1), Inches(2.4))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = w["label"] + " — " + w["title"]
    _set(r, 30, DKRED, bold=True)
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = w["subtitle"]; _set(r2, 18, GRAY, italic=True)
    mb = s.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.8))
    mr = mb.text_frame.paragraphs[0].add_run()
    mr.text = "Optimal Market Making for Cryptocurrency Options  ·  NC State University"
    _set(mr, 13, BLACK)


def _bullet_slide(prs, head, bullets, sub=None, small=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_picture(LOGO, Inches(10.7), Inches(0.3), width=Inches(2.2))
    hb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.8), Inches(0.9))
    hp = hb.text_frame; hp.word_wrap = True
    r = hp.paragraphs[0].add_run(); r.text = head; _set(r, 24, DKRED, bold=True)
    if sub:
        sr = hp.add_paragraph().add_run(); sr.text = sub; _set(sr, 13, GRAY, italic=True)
    rule = s.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(12.3), Pt(2.2))
    rule.fill.solid(); rule.fill.fore_color.rgb = RED; rule.line.fill.background()
    bb = s.shapes.add_textbox(Inches(0.7), Inches(1.75), Inches(12.0), Inches(5.4))
    btf = bb.text_frame; btf.word_wrap = True
    for i, b in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        run = p.add_run()
        sub_item = b.startswith("  ")
        run.text = ("    " + b.strip()) if sub_item else ("•  " + b)
        _set(run, 12 if small else 17, BLACK, mono=small and sub_item)
        p.space_after = Pt(5 if small else 8)


def _figure_slide(prs, w):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_picture(LOGO, Inches(10.7), Inches(0.25), width=Inches(2.2))
    hb = s.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(9.8), Inches(0.8))
    r = hb.text_frame.paragraphs[0].add_run(); r.text = "Illustration"; _set(r, 22, DKRED, bold=True)
    img = os.path.join(FIG, w["fig"])
    if os.path.exists(img):
        s.shapes.add_picture(img, Inches(1.4), Inches(1.5), width=Inches(10.5))


def _code_slide(prs, caption, code):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    hb = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.7))
    r = hb.text_frame.paragraphs[0].add_run(); r.text = "Code — " + caption; _set(r, 18, DKRED, bold=True)
    box = s.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.9))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    box.line.color.rgb = GRAY; box.line.width = Pt(0.75)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.1)
    for i, line in enumerate(code.strip("\n").split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run(); run.text = line if line else " "; _set(run, 11, BLACK, mono=True)
        p.space_after = Pt(0)


def _gloss_slide(prs, terms):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_picture(LOGO, Inches(10.7), Inches(0.3), width=Inches(2.2))
    hb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.8), Inches(0.9))
    r = hb.text_frame.paragraphs[0].add_run(); r.text = "Glossary — Key Terms"; _set(r, 24, DKRED, bold=True)
    rule = s.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(12.3), Pt(2.2))
    rule.fill.solid(); rule.fill.fore_color.rgb = RED; rule.line.fill.background()
    bb = s.shapes.add_textbox(Inches(0.7), Inches(1.75), Inches(12.0), Inches(5.4))
    tf = bb.text_frame; tf.word_wrap = True
    for i, (term, defn) in enumerate(terms):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        rt = p.add_run(); rt.text = term + " — "; _set(rt, 15, RED, bold=True)
        rd = p.add_run(); rd.text = defn; _set(rd, 15, BLACK)
        p.space_after = Pt(7)


def deck_outline(w):
    items = ["Learning objectives", w["lec1"][0], w["lec2"][0]]
    for head, _ in w.get("extra", []):
        items.append(head)
    items.append("Illustration")
    if w.get("math"):
        items.append("Derivation & rationale")
    items.append("Getting the data from Coincall")
    items.append("Code: " + w["code"][0].split(" -- ")[0])
    if w["stem"] in GLOSSARIES:
        items.append("Glossary of key terms")
    items.append("Reading")
    if w["stem"] in REFERENCES:
        items.append("References")
    return items


def build_pptx(w):
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    _title_slide(prs, w)
    _bullet_slide(prs, "Contents", deck_outline(w), small=True)
    _bullet_slide(prs, "Learning Objectives", w["objectives"])
    _bullet_slide(prs, w["lec1"][0], w["lec1"][1])
    _bullet_slide(prs, w["lec2"][0], w["lec2"][1])
    for head, bl in w.get("extra", []):
        _bullet_slide(prs, head, bl)
    _figure_slide(prs, w)
    if w.get("math"):
        _bullet_slide(prs, "Key Equations & Why We Choose Them",
            [s for eq, why in w["math"] for s in ("  " + eq, "why: " + why)][:14],
            sub="Full derivations in the PDF deck", small=True)
    _bullet_slide(prs, "Getting the Data from Coincall", w["data"],
                  sub="Cross-check exact paths at https://docs.coincall.com/", small=True)
    _code_slide(prs, w["code"][0], w["code"][1])
    if w["stem"] in GLOSSARIES:
        _gloss_slide(prs, GLOSSARIES[w["stem"]])
    _bullet_slide(prs, "Reading", w["reading"])
    if w["stem"] in REFERENCES:
        _bullet_slide(prs, "References", REFERENCES[w["stem"]], small=True)
    prs.save(os.path.join(OUT, w["stem"] + ".pptx"))


# ============================================================ BEAMER RENDER
def esc(t):
    for a, b in [("\\", r"\textbackslash{}"), ("&", r"\&"), ("%", r"\%"),
                 ("#", r"\#"), ("_", r"\_"), ("{", r"\{"), ("}", r"\}"),
                 ("$", r"\$"), ("~", r"\textasciitilde{}"), ("^", r"\textasciicircum{}"),
                 ("<", r"\textless{}"), (">", r"\textgreater{}")]:
        t = t.replace(a, b)
    return t

PRE = r"""\documentclass[aspectratio=169,10pt]{beamer}
\usetheme{default}
\setbeamertemplate{navigation symbols}{}
\usepackage{graphicx}\usepackage[T1]{fontenc}\usepackage{amsmath,amssymb}
\usepackage{listings}
\definecolor{ncsured}{HTML}{CC0000}\definecolor{ncsudk}{HTML}{990000}
\definecolor{ncsugray}{HTML}{6D6D6D}\definecolor{ncsuteal}{HTML}{008473}
\definecolor{codebg}{HTML}{F5F5F5}
\setbeamercolor{structure}{fg=ncsured}
\setbeamercolor{frametitle}{fg=white,bg=ncsured}
\setbeamercolor{title}{fg=ncsudk}
\setbeamercolor{itemize item}{fg=ncsured}
\setbeamercolor{itemize subitem}{fg=ncsured}
\setbeamerfont{frametitle}{series=\bfseries}
\setbeamertemplate{frametitle}{\nointerlineskip\vskip2pt%
  \begin{beamercolorbox}[wd=\paperwidth,ht=2.4ex,dp=1ex,leftskip=.3cm]{frametitle}%
  \small\bfseries\insertframetitle\end{beamercolorbox}}
\setbeamertemplate{footline}{\hbox{\begin{beamercolorbox}[wd=\paperwidth,ht=2.2ex,dp=1ex,leftskip=.3cm,rightskip=.3cm]{}%
  \tiny Optimal Market Making for Cryptocurrency Options \hfill __LABEL__ \hfill \insertframenumber\end{beamercolorbox}}}
\lstset{basicstyle=\ttfamily\scriptsize,backgroundcolor=\color{codebg},
  keywordstyle=\color{ncsured},commentstyle=\color{ncsugray},
  stringstyle=\color{ncsuteal},breaklines=true,columns=fullflexible,
  showstringspaces=false,frame=single,rulecolor=\color{ncsugray}}
"""

def frame(title, body):
    return r"\begin{frame}{%s}" % esc(title) + "\n" + body + "\n\\end{frame}\n"

def itemize(bullets):
    out = [r"\begin{itemize}"]
    for b in bullets:
        if b.startswith("  "):
            out.append(r"  \begin{itemize}\item \texttt{\small %s}\end{itemize}" % esc(b.strip()))
        else:
            out.append(r"  \item %s" % esc(b))
    out.append(r"\end{itemize}")
    return "\n".join(out)

def build_beamer_tex(w):
    L = [PRE.replace("__LABEL__", esc(w["label"]))]
    L.append(r"\title{%s: %s}" % (esc(w["label"]), esc(w["title"])))
    L.append(r"\subtitle{%s}" % esc(w["subtitle"]))
    L.append(r"\date{North Carolina State University \,\textbar\, Summer 2026}")
    L.append(r"\titlegraphic{\includegraphics[width=4.2cm]{%s}}" % LOGO)
    L.append(r"\begin{document}")
    L.append(r"\begin{frame}[plain]\titlepage\begin{center}\tiny " + esc(MENTOR) + r"\end{center}\end{frame}")
    # contents (table of contents)
    toc = [r"\small\begin{enumerate}\setlength{\itemsep}{3pt}"]
    for it in deck_outline(w):
        toc.append(r"  \item %s" % esc(it))
    toc.append(r"\end{enumerate}")
    L.append(r"\begin{frame}{Contents}" + "\n" + "\n".join(toc) + "\n\\end{frame}\n")
    L.append(frame("Learning Objectives", itemize(w["objectives"])))
    L.append(frame(w["lec1"][0], itemize(w["lec1"][1])))
    L.append(frame(w["lec2"][0], itemize(w["lec2"][1])))
    for head, bl in w.get("extra", []):
        L.append(frame(head, itemize(bl)))
    # illustration
    img = os.path.join(FIG, w["fig"])
    if os.path.exists(img):
        L.append(frame("Illustration",
            r"\begin{center}\includegraphics[width=0.86\textwidth]{%s}\end{center}" % img))
    # math derivation frames (with WHY)
    if w.get("math"):
        L.append(r"\section*{Derivation}")
        for eq, why in w["math"]:
            body = (r"\[%s\]" % eq) + "\n\n" + r"\small\textbf{\color{ncsured}Why:} " + why
            L.append(r"\begin{frame}{Derivation \& Rationale}" + "\n" + body + "\n\\end{frame}\n")
    # data
    L.append(frame("Getting the Data from Coincall", itemize(w["data"]) +
        "\n\\vspace{2pt}\\tiny\\textit{Cross-check exact paths at https://docs.coincall.com/}"))
    # code (fragile frame with listings)
    L.append(r"\begin{frame}[fragile]{Code --- %s}" % esc(w["code"][0].split(" -- ")[0]))
    L.append(r"\begin{lstlisting}" + "\n" + w["code"][1].strip("\n") + "\n" + r"\end{lstlisting}")
    L.append(r"\end{frame}")
    # glossary
    if w["stem"] in GLOSSARIES:
        gl = [r"\footnotesize\begin{description}"]
        for term, defn in GLOSSARIES[w["stem"]]:
            gl.append(r"\item[\color{ncsured}%s] %s" % (esc(term), esc(defn)))
        gl.append(r"\end{description}")
        L.append(r"\begin{frame}{Glossary --- Key Terms}" + "\n"
                 + "\n".join(gl) + "\n\\end{frame}\n")
    L.append(frame("Reading", itemize(w["reading"])))
    # references (full citations)
    if w["stem"] in REFERENCES:
        rf = [r"\scriptsize\begin{itemize}\setlength{\itemsep}{3pt}"]
        for c in REFERENCES[w["stem"]]:
            rf.append(r"  \item %s" % esc(c))
        rf.append(r"\end{itemize}")
        L.append(r"\begin{frame}{References}" + "\n" + "\n".join(rf) + "\n\\end{frame}\n")
    L.append(r"\end{document}")
    open(os.path.join(OUT, w["stem"] + ".tex"), "w", encoding="utf-8").write("\n".join(L))


if __name__ == "__main__":
    for w in WEEKS:
        build_pptx(w); build_beamer_tex(w)
        print("built", w["stem"])
    print("Done: %d weeks." % len(WEEKS))
